"""
Image extraction from PDF/PPTX/DOCX source documents.

Pure Python implementation — no Rust commands required.
Uses:
  - PyMuPDF (fitz) for PDF images
  - python-pptx for PPTX slides
  - python-docx for DOCX paragraphs
  - zipfile as fallback for raw ZIP-based extraction

Extracted images are saved to:
  `<project>/wiki/media/<slug>/`

Reference: nashsu/llm_wiki/src/lib/extract-source-images.ts
"""

from __future__ import annotations

import hashlib
import os
import shutil
import zipfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional


# ── Types ──────────────────────────────────────────────────────────────

SUPPORTED_PDF_EXTS = ("pdf",)
SUPPORTED_OFFICE_EXTS = ("pptx", "docx", "ppt", "doc")


@dataclass
class SavedImage:
    index: int
    mime_type: str
    page: Optional[int]   # PDF page (1-based) or None for Office
    width: int
    height: int
    rel_path: str         # relative to wiki/ root, e.g. "media/slug/img-1.png"
    abs_path: str         # absolute filesystem path
    sha256: str


# ── Main dispatcher ───────────────────────────────────────────────────

def extract_and_save_source_images(
    project_path: str,
    source_path: str,
) -> list[SavedImage]:
    """Extract every embedded image from `source_path` and save to wiki/media/<slug>/.

    Returns a list of SavedImage metadata. Returns [] for unsupported file types
    or when the source has no extractable images. Errors are logged and returned
    as empty — image extraction failure must NEVER abort the ingest pipeline.
    """
    ext = Path(source_path).suffix.lower().lstrip(".")

    if ext in SUPPORTED_PDF_EXTS:
        return _extract_pdf_images(project_path, source_path)
    elif ext in SUPPORTED_OFFICE_EXTS:
        return _extract_office_images(project_path, source_path)
    else:
        return []


# ── PDF extraction ─────────────────────────────────────────────────────

def _extract_pdf_images(project_path: str, source_path: str) -> list[SavedImage]:
    """Extract images from PDF using PyMuPDF."""
    try:
        import fitz  # PyMuPDF
    except ImportError:
        print("[extract-images] PyMuPDF (fitz) not installed — PDF images skipped")
        return []

    slug = Path(source_path).stem
    dest_dir = Path(project_path) / "wiki" / "media" / slug
    rel_base = f"media/{slug}"

    images: list[SavedImage] = []
    img_index = 0

    try:
        doc = fitz.open(source_path)
    except Exception as e:
        print(f"[extract-images] failed to open PDF {source_path}: {e}")
        return []

    for page_num in range(len(doc)):
        page = doc[page_num]
        image_list = page.get_images(full=True)
        for img in image_list:
            xref = img[0]
            try:
                base_image = doc.extract_image(xref)
                img_bytes = base_image["image"]
                mime = base_image["ext"]
                width = base_image["width"]
                height = base_image["height"]
            except Exception:
                continue

            ext = _mime_to_ext(mime)
            filename = f"img-{img_index + 1}.{ext}"
            abs_path = dest_dir / filename
            rel_path = f"{rel_base}/{filename}"

            # Save image bytes
            try:
                dest_dir.mkdir(parents=True, exist_ok=True)
                abs_path.write_bytes(img_bytes)
            except OSError as e:
                print(f"[extract-images] failed to write {abs_path}: {e}")
                continue

            sha = hashlib.sha256(img_bytes).hexdigest()
            images.append(SavedImage(
                index=img_index,
                mime_type=mime,
                page=page_num + 1,
                width=width,
                height=height,
                rel_path=rel_path,
                abs_path=str(abs_path),
                sha256=sha,
            ))
            img_index += 1

    doc.close()
    return images


# ── Office extraction ─────────────────────────────────────────────────

def _extract_office_images(project_path: str, source_path: str) -> list[SavedImage]:
    """Extract images from PPTX/DOCX using python-pptx / python-docx or zipfile fallback."""
    ext = Path(source_path).suffix.lower().lstrip(".")

    # Try python-pptx / python-docx first
    if ext in ("pptx",):
        try:
            return _extract_pptx_images(project_path, source_path)
        except Exception as e:
            print(f"[extract-images] python-pptx failed for {source_path}: {e}")
    elif ext in ("docx",):
        try:
            return _extract_docx_images(project_path, source_path)
        except Exception as e:
            print(f"[extract-images] python-docx failed for {source_path}: {e}")

    # Fallback: raw ZIP extraction
    return _extract_office_zip_fallback(project_path, source_path)


def _extract_pptx_images(project_path: str, source_path: str) -> list[SavedImage]:
    """Extract images from PPTX using python-pptx."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    slug = Path(source_path).stem
    dest_dir = Path(project_path) / "wiki" / "media" / slug
    rel_base = f"media/{slug}"

    images: list[SavedImage] = []
    img_index = 0

    prs = Presentation(source_path)
    for slide_num, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img = shape.image
                img_bytes = img.blob
                mime = img.content_type
                ext = _mime_to_ext(mime)
                filename = f"img-{img_index + 1}.{ext}"
                abs_path = dest_dir / filename
                rel_path = f"{rel_base}/{filename}"

                try:
                    dest_dir.mkdir(parents=True, exist_ok=True)
                    abs_path.write_bytes(img_bytes)
                except OSError:
                    continue

                sha = hashlib.sha256(img_bytes).hexdigest()
                images.append(SavedImage(
                    index=img_index,
                    mime_type=mime,
                    page=slide_num,
                    width=img.width,
                    height=img.height,
                    rel_path=rel_path,
                    abs_path=str(abs_path),
                    sha256=sha,
                ))
                img_index += 1

    return images


def _extract_docx_images(project_path: str, source_path: str) -> list[SavedImage]:
    """Extract images from DOCX using python-docx."""
    from docx import Document

    slug = Path(source_path).stem
    dest_dir = Path(project_path) / "wiki" / "media" / slug
    rel_base = f"media/{slug}"

    images: list[SavedImage] = []
    img_index = 0

    doc = Document(source_path)
    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            try:
                img_bytes = rel.target_part.blob
                mime = rel.target_part.content_type
                ext = _mime_to_ext(mime)
                filename = f"img-{img_index + 1}.{ext}"
                abs_path = dest_dir / filename
                rel_path = f"{rel_base}/{filename}"

                dest_dir.mkdir(parents=True, exist_ok=True)
                abs_path.write_bytes(img_bytes)

                sha = hashlib.sha256(img_bytes).hexdigest()
                images.append(SavedImage(
                    index=img_index,
                    mime_type=mime,
                    page=None,
                    width=0,
                    height=0,
                    rel_path=rel_path,
                    abs_path=str(abs_path),
                    sha256=sha,
                ))
                img_index += 1
            except Exception:
                continue

    return images


def _extract_office_zip_fallback(project_path: str, source_path: str) -> list[SavedImage]:
    """Fallback: extract images from PPTX/DOCX using raw zipfile."""
    slug = Path(source_path).stem
    dest_dir = Path(project_path) / "wiki" / "media" / slug
    rel_base = f"media/{slug}"

    images: list[SavedImage] = []
    img_index = 0

    try:
        with zipfile.ZipFile(source_path, "r") as zf:
            for name in zf.namelist():
                if not name.startswith("ppt/media/") and not name.startswith("word/media/"):
                    continue
                if not (_is_image_ext(name)):
                    continue

                try:
                    img_bytes = zf.read(name)
                except zipfile.BadZipFile:
                    continue

                ext = Path(name).suffix.lower().lstrip(".") or "png"
                mime = _ext_to_mime(ext)
                filename = f"img-{img_index + 1}.{ext}"
                abs_path = dest_dir / filename
                rel_path = f"{rel_base}/{filename}"

                try:
                    dest_dir.mkdir(parents=True, exist_ok=True)
                    abs_path.write_bytes(img_bytes)
                except OSError:
                    continue

                sha = hashlib.sha256(img_bytes).hexdigest()
                images.append(SavedImage(
                    index=img_index,
                    mime_type=mime,
                    page=None,
                    width=0,
                    height=0,
                    rel_path=rel_path,
                    abs_path=str(abs_path),
                    sha256=sha,
                ))
                img_index += 1
    except (zipfile.BadZipFile, OSError) as e:
        print(f"[extract-images] zipfile fallback failed for {source_path}: {e}")

    return images


# ── Markdown section builder ───────────────────────────────────────────

def build_image_markdown_section(
    images: list[SavedImage],
    captions_by_sha: Optional[dict[str, str]] = None,
) -> str:
    """Build the markdown section to splice into source content for the LLM.

    Each image is referenced by its rel_path with a caption (or empty alt
    text if captions aren't available yet).

    Returns an empty string when there are no images.
    """
    if not images:
        return ""

    captions = captions_by_sha or {}

    # Group by page
    by_page: dict[str, list[SavedImage]] = {}
    for img in images:
        key = f"Page {img.page}" if img.page is not None else "Document"
        by_page.setdefault(key, []).append(img)

    # Sort page keys numerically, "Document" last
    def page_sort_key(key: str) -> tuple[int, int]:
        if key == "Document":
            return (99999, 0)
        num = int("".join(filter(str.isdigit, key))) or 0
        return (0, num)

    ordered_keys = sorted(by_page.keys(), key=page_sort_key)

    def sanitize(s: str) -> str:
        return s.replace("\r\n", " ").replace("\n", " ").replace("]", ")").strip()

    lines: list[str] = ["", "", "## Embedded Images", ""]
    for key in ordered_keys:
        lines.append(f"### {key}")
        lines.append("")
        for img in by_page[key]:
            caption = captions.get(img.sha256, "")
            alt = sanitize(caption) if caption else ""
            lines.append(f"![{alt}]({img.rel_path})")
        lines.append("")

    return "\n".join(lines)


# ── Helpers ────────────────────────────────────────────────────────────

_MIME_TO_EXT = {
    "image/png": "png",
    "image/jpeg": "jpg",
    "image/jpg": "jpg",
    "image/gif": "gif",
    "image/webp": "webp",
    "image/bmp": "bmp",
    "image/svg+xml": "svg",
}

_EXT_TO_MIME = {v: k for k, v in _MIME_TO_EXT.items()}


def _mime_to_ext(mime: str) -> str:
    return _MIME_TO_EXT.get(mime, "png")


def _ext_to_mime(ext: str) -> str:
    return _EXT_TO_MIME.get(ext.lower(), "image/png")


_IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".svg", ".tiff"}


def _is_image_ext(name: str) -> bool:
    return Path(name).suffix.lower() in _IMAGE_EXTS
